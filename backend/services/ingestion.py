from pathlib import Path
from typing import Iterable

from fastapi import UploadFile
from sqlalchemy.orm import Session

from backend.models import Document, DocumentChunk, DocumentStatus
from backend.settings import settings
from backend.services.storage import save_upload_file
from backend.services.indexing import index_chunks, mark_all_docs_outdated


try:
    from pypdf import PdfReader
except ImportError:  # pragma: no cover
    PdfReader = None

try:
    from docx import Document as DocxDocument
except ImportError:  # pragma: no cover
    DocxDocument = None

try:
    from openpyxl import load_workbook
except ImportError:  # pragma: no cover
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None


def extract_pages(path: Path) -> list[dict]:
    ext = path.suffix.lower()
    pages: list[dict] = []

    if ext == ".pdf" and PdfReader is not None:
        reader = PdfReader(str(path))
        for idx, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append({"text": text, "page": idx + 1, "bbox": None})
        return pages

    if ext == ".docx" and DocxDocument is not None:
        doc = DocxDocument(str(path))
        text = "\n".join(p.text for p in doc.paragraphs if p.text)
        pages.append({"text": text, "page": None, "bbox": None})
        return pages

    if ext == ".xlsx" and load_workbook is not None:
        wb = load_workbook(str(path), data_only=True)
        for sheet in wb.worksheets:
            lines = []
            for row in sheet.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    lines.append(row_text)
            pages.append({"text": "\n".join(lines), "page": None, "bbox": None})
        return pages

    if ext == ".pptx" and Presentation is not None:
        pres = Presentation(str(path))
        for slide_idx, slide in enumerate(pres.slides, start=1):
            lines = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text:
                        lines.append(shape.text)
            pages.append({"text": "\n".join(lines), "page": slide_idx, "bbox": None})
        return pages

    text = path.read_text(encoding="utf-8", errors="ignore")
    pages.append({"text": text, "page": None, "bbox": None})
    return pages


def chunk_pages(pages: Iterable[dict], chunk_size: int, overlap: int) -> list[dict]:
    chunks: list[dict] = []
    for page in pages:
        text = page["text"] or ""
        if not text.strip():
            continue
        start = 0
        while start < len(text):
            end = min(start + chunk_size, len(text))
            chunk_text = text[start:end]
            if chunk_text.strip():
                chunks.append({
                    "text": chunk_text,
                    "page": page["page"],
                    "bbox": page["bbox"],
                })
            if end == len(text):
                break
            start = end - overlap
    return chunks


def process_document(db: Session, doc: Document) -> Document:
    dest = Path(doc.storage_path)
    pages = extract_pages(dest)
    raw_chunks = chunk_pages(pages, settings.chunk_size, settings.chunk_overlap)

    chunk_models: list[DocumentChunk] = []
    for idx, chunk in enumerate(raw_chunks):
        chunk_models.append(DocumentChunk(
            document_id=doc.id,
            chunk_index=idx,
            text=chunk["text"],
            page=chunk["page"],
            bbox=chunk["bbox"],
        ))
    db.add_all(chunk_models)
    doc.status = DocumentStatus.PARSED
    db.commit()

    index_chunks(db, doc, chunk_models)
    doc.status = DocumentStatus.INDEXED
    db.commit()

    mark_all_docs_outdated(db)
    return doc


def ingest_document(db: Session, upload: UploadFile) -> Document:
    doc = Document(
        filename=upload.filename,
        content_type=upload.content_type or "application/octet-stream",
        status=DocumentStatus.UPLOADED,
        storage_path="",
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)

    dest = Path(settings.storage_path) / f"{doc.id}_{upload.filename}"
    save_upload_file(upload, dest)
    doc.storage_path = str(dest)
    db.commit()

    return process_document(db, doc)
